from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Helper function to add a slide with title and content
    def add_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0

    # 1. Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Twitter Sentiment Analyzer"
    subtitle.text = "From Unsupervised Clustering to Supervised Classification\nProject Documentation & Architecture"

    # 2. Project Overview
    add_slide("Project Overview", [
        "Goal: Categorize tweets into Positive, Negative, Neutral, and Irrelevant classes.",
        "Challenge: Handling noisy social media data (URLs, hashtags, @mentions).",
        "Solution: End-to-end Machine Learning pipeline with a Streamlit Web UI."
    ])

    # 3. Project Structure
    add_slide("Modular Project Structure", [
        "src/components: Data Ingestion, Transformation, and Model Training logic.",
        "src/pipeline: Automated Train and Predict workflows.",
        "artifacts/: Persistent storage for Model (.pkl) and Vectorizer.",
        "logs/: Comprehensive logging for debugging and tracking experiments."
    ])

    # 4. Why NOT Unsupervised KMeans?
    add_slide("The Case Against Unsupervised KMeans", [
        "Cluster Limitation: Groups data by word frequency, not emotional intent.",
        "Topic Bias: Clusters often group by subject (e.g., 'Apple' vs 'Microsoft') instead of sentiment.",
        "Low Accuracy: Achieved ~30% accuracy (barely better than random guessing).",
        "Label Guessing: Requires manual 'majority vote' to assign names to clusters."
    ])

    # 5. Why Logistic Regression?
    add_slide("The Power of Logistic Regression", [
        "Supervised Learning: Uses ground-truth labels to learn specific patterns.",
        "Weighted Features: Assigns numerical importance to words (e.g., 'love' = +2.5, 'fail' = -3.0).",
        "Performance: Jumped from 30% to 80%+ accuracy instantly.",
        "Speed: Extremely efficient for high-dimensional text data."
    ])

    # 6. NLP Transformation Pipeline
    add_slide("The NLP Engine", [
        "Cleaning: Regex-based removal of noise (URLs, Punctuation, Digits).",
        "Lemmatization: Reducing words to their dictionary root (WordNet).",
        "TF-IDF: Term Frequency-Inverse Document Frequency for feature extraction.",
        "N-grams: Capturing context using Unigrams and Bigrams (e.g., 'not good')."
    ])

    # 7. Web Interface (Streamlit)
    add_slide("Interactive Web UI", [
        "Single Prediction: Real-time analysis for instant feedback.",
        "Batch Processing: CSV upload for large-scale data analysis.",
        "Visualizations: Plotly-powered distribution charts and count plots.",
        "User Experience: Responsive design with sentiment-themed emojis and colors."
    ])

    # 8. Conclusion
    add_slide("Conclusion & Next Steps", [
        "Success: Validated the superiority of supervised learning for this task.",
        "Key takeaway: Modular code makes the pipeline easy to maintain and scale.",
        "Future: Potential migration to BERT or Transformers for deep contextual analysis."
    ])

    # Save the presentation
    file_name = "Twitter_Sentiment_Analysis.pptx"
    prs.save(file_name)
    print(f"Presentation saved as {file_name}")

if __name__ == "__main__":
    create_presentation()